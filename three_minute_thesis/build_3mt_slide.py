#!/usr/bin/env python3
"""
Build a single static 3MT slide for the B.A.I.L.I.F.F. presentation.

Output:
  three_minute_thesis/output/BAILIFF_3MT_slide.pptx
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
OUT_DIR = ROOT / "output"
OUT_PPTX = OUT_DIR / "BAILIFF_3MT_slide.pptx"
CHART = ROOT / "output" / "3mt_flip_rate_chart.png"
FALLBACK_CHART = ROOT.parent / "paper" / "plots" / "family_effect_snapshot.png"


def add_textbox(slide, left, top, width, height, text, font_size=20, bold=False, color=RGBColor(20, 20, 20), align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def main():
    OUT_DIR.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Background panel
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(248, 250, 252)
    bg.line.fill.background()

    # Title and subtitle
    add_textbox(
        slide,
        Inches(0.5),
        Inches(0.25),
        Inches(12.3),
        Inches(0.6),
        "Can AI Hold a Fair Trial?",
        font_size=38,
        bold=True,
        color=RGBColor(15, 23, 42),
        align=PP_ALIGN.LEFT,
    )
    add_textbox(
        slide,
        Inches(0.5),
        Inches(0.9),
        Inches(12.3),
        Inches(0.4),
        "B.A.I.L.I.F.F. audits legal AI with adversarial, paired mock trials",
        font_size=18,
        color=RGBColor(51, 65, 85),
    )

    # Left content card
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.45), Inches(5.6), Inches(4.9))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(255, 255, 255)
    card.line.color.rgb = RGBColor(203, 213, 225)
    card.line.width = Pt(1.2)

    add_textbox(
        slide,
        Inches(0.8),
        Inches(1.75),
        Inches(5.0),
        Inches(0.35),
        "What we tested",
        font_size=20,
        bold=True,
        color=RGBColor(30, 41, 59),
    )

    bullet_text = (
        "• Same case facts, only demographic cue toggled\n"
        "• Three-agent courtroom: Judge, Prosecution, Defense\n"
        "• We measured both:\n"
        "   - Verdict direction (delta, OR)\n"
        "   - Procedural stability (flip rate)"
    )
    box = add_textbox(
        slide,
        Inches(0.8),
        Inches(2.15),
        Inches(5.0),
        Inches(2.2),
        bullet_text,
        font_size=16,
        color=RGBColor(51, 65, 85),
    )
    box.text_frame.word_wrap = True

    # Key message banner
    banner = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.55), Inches(5.0), Inches(1.55))
    banner.fill.solid()
    banner.fill.fore_color.rgb = RGBColor(15, 23, 42)
    banner.line.fill.background()
    key_text = (
        "Key message:\n"
        "Even when pooled outcome shifts are modest,\n"
        "flip rates stay non-zero across model families."
    )
    key_box = add_textbox(
        slide,
        Inches(1.05),
        Inches(4.8),
        Inches(4.55),
        Inches(1.2),
        key_text,
        font_size=14,
        bold=False,
        color=RGBColor(248, 250, 252),
    )
    key_box.text_frame.word_wrap = True

    # Right chart panel
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.3), Inches(1.45), Inches(6.5), Inches(4.9))
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(255, 255, 255)
    panel.line.color.rgb = RGBColor(203, 213, 225)
    panel.line.width = Pt(1.2)

    add_textbox(
        slide,
        Inches(6.6),
        Inches(1.75),
        Inches(5.9),
        Inches(0.35),
        "Observed stability by model family",
        font_size=20,
        bold=True,
        color=RGBColor(30, 41, 59),
    )

    chart_path = CHART if CHART.exists() else FALLBACK_CHART
    if chart_path.exists():
        slide.shapes.add_picture(str(chart_path), Inches(6.55), Inches(2.15), width=Inches(6.0), height=Inches(3.95))
    else:
        add_textbox(
            slide,
            Inches(6.8),
            Inches(3.2),
            Inches(5.4),
            Inches(1.0),
            "Chart not found:\nthree_minute_thesis/output/3mt_flip_rate_chart.png",
            font_size=15,
            color=RGBColor(185, 28, 28),
            align=PP_ALIGN.CENTER,
        )

    # Footer
    add_textbox(
        slide,
        Inches(0.5),
        Inches(6.95),
        Inches(12.4),
        Inches(0.3),
        "Single static slide (3MT-compliant): no animations, no transitions, no embedded media",
        font_size=11,
        color=RGBColor(71, 85, 105),
        align=PP_ALIGN.CENTER,
    )

    prs.save(str(OUT_PPTX))
    print(f"Wrote {OUT_PPTX}")


if __name__ == "__main__":
    main()
